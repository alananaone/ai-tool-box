# workflow_scripts/summary_to_ppt.py (再次優化版)
import os
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR # 新增 MSO_VERTICAL_ANCHOR
import traceback
import logging

# 字型設定 - 全部改為標楷體
FONT_PRIMARY = '標楷體' # 主要字型，用於所有元素

# 投影片版面配置索引 (基於標準 Office 佈景主題 - 索引可能因母片不同而異)
LAYOUT_H1_TITLE_ONLY = 5
LAYOUT_H2_TITLE_AND_CONTENT = 1

def get_indent_level(para):
    """計算 Word 段落的縮排層級，用於 PPT 的項目符號層級。"""
    indent_val = 0
    base_indent_inches = 0.1
    indent_step_inches = 0.4

    if para.paragraph_format.left_indent:
        try:
            indent_val = para.paragraph_format.left_indent.inches
        except AttributeError:
            indent_val = 0
        except Exception:
            indent_val = 0
            logging.debug(f"無法獲取段落縮進的 .inches 值: '{para.text[:30]}...'")
    level = max(0, int(round((indent_val - base_indent_inches) / indent_step_inches)))
    return min(level, 8)

def run_conversion_to_ppt(input_summary_path: str, output_ppt_path: str) -> bool:
    try:
        logging.info(f"  開始轉換 Word 摘要 '{os.path.basename(input_summary_path)}' 到 PPTX...")
        doc = DocxDocument(input_summary_path)
        prs = Presentation()

        current_h2_slide = None
        current_h2_content_placeholder = None
        current_h2_title_for_slide = ""
        content_item_count_on_current_slide = 0
        MAX_ITEMS_PER_SLIDE = 7

        for para_idx, para in enumerate(doc.paragraphs):
            stripped_para_text = para.text.strip()
            if not stripped_para_text:
                logging.debug(f"    跳過空段落 (Word 段落索引 {para_idx})")
                continue
            style_name = para.style.name
            logging.debug(f"    處理 Word 段落 {para_idx}: '{stripped_para_text[:50]}', 樣式: '{style_name}'")

            is_h1 = style_name.startswith('Heading 1')
            is_h2 = style_name.startswith('Heading 2')

            if is_h1:
                try:
                    slide_layout = prs.slide_layouts[LAYOUT_H1_TITLE_ONLY]
                    h1_slide = prs.slides.add_slide(slide_layout)
                    
                    title_shape_h1 = None
                    if getattr(h1_slide.shapes, 'title', None) is not None:
                        title_shape_h1 = h1_slide.shapes.title
                    elif h1_slide.placeholders and len(h1_slide.placeholders) > 0:
                        title_shape_h1 = h1_slide.placeholders[0]

                    if title_shape_h1:
                        title_shape_h1.text = stripped_para_text
                        if title_shape_h1.has_text_frame:
                            # 設定文字框內垂直置中
                            title_shape_h1.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                            if title_shape_h1.text_frame.paragraphs:
                                tf = title_shape_h1.text_frame.paragraphs[0]
                                tf.font.name = FONT_PRIMARY # 使用統一字型
                                tf.font.size = Pt(40)
                                tf.font.bold = True
                                # 水平置中 (如果版面本身不是置中對齊)
                                tf.alignment = PP_ALIGN.CENTER 
                        logging.info(f"    已建立 H1 章節頁: '{stripped_para_text[:50]}...'")
                    else:
                        logging.warning(f"    H1 章節頁 ('{stripped_para_text[:50]}...') 的版面配置 (索引 {LAYOUT_H1_TITLE_ONLY}) "
                                        f"可能沒有預期的標題佔位符。")

                    current_h2_slide = None
                    current_h2_content_placeholder = None
                    content_item_count_on_current_slide = 0
                    current_h2_title_for_slide = ""

                except IndexError:
                    logging.error(f"    錯誤：找不到索引為 {LAYOUT_H1_TITLE_ONLY} 的 H1 ('僅標題') 投影片版面配置。跳過 H1 頁面 '{stripped_para_text[:50]}...' 的創建。")
                except Exception as e_h1_slide:
                    logging.error(f"    建立 H1 章節頁 '{stripped_para_text[:50]}...' 時發生錯誤: {e_h1_slide}", exc_info=True)

            elif is_h2:
                try:
                    slide_layout = prs.slide_layouts[LAYOUT_H2_TITLE_AND_CONTENT]
                    current_h2_slide = prs.slides.add_slide(slide_layout)
                    current_h2_title_for_slide = stripped_para_text

                    title_shape_h2 = current_h2_slide.shapes.title
                    if title_shape_h2:
                        title_shape_h2.text = current_h2_title_for_slide
                        if title_shape_h2.has_text_frame and title_shape_h2.text_frame.paragraphs:
                            tf = title_shape_h2.text_frame.paragraphs[0]
                            tf.font.name = FONT_PRIMARY # 使用統一字型
                            tf.font.size = Pt(32)
                            tf.font.bold = True
                    else:
                        logging.warning(f"    H2 內容頁 ('{current_h2_title_for_slide[:50]}...') 的版面配置 (索引 {LAYOUT_H2_TITLE_AND_CONTENT}) "
                                        f"缺少 'title' shape。")

                    if len(current_h2_slide.placeholders) > 1 and current_h2_slide.placeholders[1].has_text_frame:
                        current_h2_content_placeholder = current_h2_slide.placeholders[1]
                        current_h2_content_placeholder.text_frame.clear()
                        # 設定內容佔位符文字框的預設字型 (這樣裡面的每個段落預設就是這個字型)
                        # 雖然下面還是會為每個段落設定，但這裡設定可以作為一個基礎
                        current_h2_content_placeholder.text_frame.paragraphs[0].font.name = FONT_PRIMARY

                        if not current_h2_content_placeholder.text_frame.paragraphs: # clear 後可能為空
                             p_temp = current_h2_content_placeholder.text_frame.add_paragraph()
                             p_temp.font.name = FONT_PRIMARY # 確保新段落也有字型

                        content_item_count_on_current_slide = 0
                        logging.info(f"    已建立 H2 內容頁: '{current_h2_title_for_slide[:50]}...'")
                    else:
                        logging.error(f"    錯誤：在 H2 版面配置 {LAYOUT_H2_TITLE_AND_CONTENT} 中找不到預期的內容佔位符 (placeholders[1]) "
                                      f"或其沒有 text_frame。")
                        current_h2_content_placeholder = None
                
                except IndexError:
                    logging.error(f"    錯誤：找不到索引為 {LAYOUT_H2_TITLE_AND_CONTENT} 的 H2 ('標題及內容') 投影片版面配置。")
                    current_h2_slide = None; current_h2_content_placeholder = None
                except Exception as e_h2_slide:
                    logging.error(f"    建立 H2 內容頁 '{stripped_para_text[:50]}...' 時發生錯誤: {e_h2_slide}", exc_info=True)
                    current_h2_slide = None; current_h2_content_placeholder = None

            elif current_h2_slide and current_h2_content_placeholder:
                if content_item_count_on_current_slide >= MAX_ITEMS_PER_SLIDE:
                    logging.info(f"    項目數達上限 ({MAX_ITEMS_PER_SLIDE})，為 H2 '{current_h2_title_for_slide[:50]}' 建立接續頁...")
                    try:
                         slide_layout = prs.slide_layouts[LAYOUT_H2_TITLE_AND_CONTENT]
                         current_h2_slide = prs.slides.add_slide(slide_layout)
                         
                         title_shape_cont = current_h2_slide.shapes.title
                         if title_shape_cont:
                             title_shape_cont.text = f"{current_h2_title_for_slide} (續)"
                             if title_shape_cont.has_text_frame and title_shape_cont.text_frame.paragraphs:
                                tf = title_shape_cont.text_frame.paragraphs[0]
                                tf.font.name = FONT_PRIMARY; tf.font.size = Pt(32); tf.font.bold = True
                         
                         if len(current_h2_slide.placeholders) > 1 and current_h2_slide.placeholders[1].has_text_frame:
                            current_h2_content_placeholder = current_h2_slide.placeholders[1]
                            current_h2_content_placeholder.text_frame.clear()
                            current_h2_content_placeholder.text_frame.paragraphs[0].font.name = FONT_PRIMARY # 設定基礎字型
                            if not current_h2_content_placeholder.text_frame.paragraphs:
                                p_temp = current_h2_content_placeholder.text_frame.add_paragraph()
                                p_temp.font.name = FONT_PRIMARY

                            content_item_count_on_current_slide = 0
                            logging.info(f"    已建立 H2 '{current_h2_title_for_slide[:50]}' 的接續頁。")
                         else:
                            current_h2_content_placeholder = None
                            logging.error("    錯誤：H2 接續頁找不到有效的內容佔位符。")
                    except IndexError:
                         current_h2_slide = None; current_h2_content_placeholder = None
                         logging.error("    錯誤：無法建立 H2 內容的接續頁面 (找不到版面配置)。")
                    except Exception as e_cont_slide:
                        current_h2_slide = None; current_h2_content_placeholder = None
                        logging.error(f"    建立 H2 內容接續頁面時發生錯誤: {e_cont_slide}", exc_info=True)

                if current_h2_content_placeholder:
                    p = current_h2_content_placeholder.text_frame.add_paragraph()
                    p.text = stripped_para_text
                    p.level = get_indent_level(para)
                    p.font.name = FONT_PRIMARY # 使用統一字型
                    p.font.size = Pt(20)
                    p.font.bold = False
                    content_item_count_on_current_slide += 1
                    logging.debug(f"      已添加內容到 H2 '{current_h2_title_for_slide[:30]}...': '{stripped_para_text[:30]}', 層級: {p.level}")
            
            elif not is_h1 and not is_h2 and para_idx == 0:
                 logging.warning(f"    文件開頭段落 '{stripped_para_text[:50]}...' 不是 H1 或 H2，將被忽略。")
            elif not is_h1 and not is_h2 and not current_h2_slide:
                 logging.debug(f"    段落 '{stripped_para_text[:50]}...' (非H1/H2) 出現在 H1 之後但 H2 之前，或無任何有效 H2 頁面，將被忽略。")


        if not prs.slides:
             logging.warning(f"警告：文件 '{os.path.basename(input_summary_path)}' 未能生成任何投影片。請檢查 Word 文件是否包含有效的 H1/H2 結構。")
             return False
        else:
             prs.save(output_ppt_path)
             logging.info(f"  PPTX 簡報 '{os.path.basename(output_ppt_path)}' 儲存成功。共產生 {len(prs.slides)} 張投影片。")
             return True

    except Exception as e:
        logging.error(f"!!!!!!!!!! 處理 Word 檔案 '{os.path.basename(input_summary_path)}' 轉換為 PPTX 時發生嚴重錯誤 !!!!!!!!!!")
        logging.error(traceback.format_exc())
        return False